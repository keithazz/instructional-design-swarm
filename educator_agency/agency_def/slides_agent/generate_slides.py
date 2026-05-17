"""Tool: generate educator-style slides from structured slide briefs.

Per `.claude/plans/i-want-to-start-stateless-oasis.md` §6.8 and the
`002_presentation_styling` follow-up plan: the slides agent emits one
`SlideBrief` per slide; this tool calls a stateless HTML-writer sub-agent
per slide (driven by `html_writer_instructions.md` in this directory),
concatenates the resulting HTML fragments, wraps each in the course
style.css, runs `html2pptx_runner.js`, then writes the PPTX through the
active `FileOpsBackend`.

Architectural anchors:
- The tool's external contract stays one-shot and backend-gated.
- Sub-agents are internal, stateless, and never write to disk.
- All file writes route through `get_backend()`.
"""

from __future__ import annotations

import logging
import os
import re
import subprocess
import sys
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Any

from agency_swarm import Agent, ModelSettings, Reasoning
from agency_swarm.tools import BaseTool
from agents.extensions.models.litellm_model import LitellmModel
from openai import AsyncOpenAI
from pydantic import BaseModel, Field

from educator_agency.runtime._context import get_backend
from educator_agency.runtime.file_ops import Accepted, Failed, Pending, Rejected

RUNNER_JS = Path(__file__).parent.parent.parent.parent / "slides_agent" / "tools" / "html2pptx_runner.js"
_HTML_WRITER_INSTRUCTIONS = Path(__file__).with_name("html_writer_instructions.md")
_HTML_WRITER_MODEL_CLAUDE = "anthropic/claude-sonnet-4-6"
_HTML_WRITER_MODEL_OAI = "gpt-5.3-codex"

# ---------------------------------------------------------------------------
# Logging
# ---------------------------------------------------------------------------
# Logger is configured eagerly at import time so events flow regardless of
# how the host process set up logging. Two sinks:
#   1. stderr at INFO  — per-slide events visible in the TUI/server console.
#   2. <debug_dir>/slides.log at DEBUG — full prompts, raw responses, etc.
# Per-slide artifacts (prompt + raw response + extracted fragment) are also
# dumped to <debug_dir>/<run_timestamp>/slide_NNN.* so they can be diffed.
#
# debug_dir resolution:
#   $SLIDES_DEBUG_DIR if set, else Path.cwd() / ".slides_debug"
# To silence: set $SLIDES_LOG=quiet (drops the stderr handler).
# To go fully verbose on stderr too: set $SLIDES_LOG=debug.

logger = logging.getLogger("educator_agency.slides")


def _debug_dir() -> Path:
    base = os.getenv("SLIDES_DEBUG_DIR")
    return Path(base).expanduser() if base else Path.cwd() / ".slides_debug"


def _configure_logging() -> None:
    """Idempotent: only attaches handlers the first time it's called."""
    if getattr(logger, "_slides_configured", False):
        return
    logger.setLevel(logging.DEBUG)
    logger.propagate = False  # avoid double-logging via root

    fmt = logging.Formatter(
        "%(asctime)s [%(levelname)s] %(message)s", datefmt="%H:%M:%S"
    )

    stderr_mode = os.getenv("SLIDES_LOG", "info").lower()
    if stderr_mode != "quiet":
        stderr_handler = logging.StreamHandler(sys.stderr)
        stderr_handler.setLevel(
            logging.DEBUG if stderr_mode == "debug" else logging.INFO
        )
        stderr_handler.setFormatter(fmt)
        logger.addHandler(stderr_handler)

    try:
        debug_root = _debug_dir()
        debug_root.mkdir(parents=True, exist_ok=True)
        file_handler = logging.FileHandler(debug_root / "slides.log", mode="a", encoding="utf-8")
        file_handler.setLevel(logging.DEBUG)
        file_handler.setFormatter(fmt)
        logger.addHandler(file_handler)
    except OSError as exc:  # debug dir is optional; don't crash the tool
        logger.warning("Could not attach slides.log file handler: %s", exc)

    logger._slides_configured = True  # type: ignore[attr-defined]


_configure_logging()


def _make_run_artifact_dir() -> Path | None:
    """Create a per-run subdir under the debug dir, or return None on failure."""
    try:
        run_dir = _debug_dir() / datetime.now().strftime("%Y%m%d-%H%M%S")
        run_dir.mkdir(parents=True, exist_ok=True)
        return run_dir
    except OSError as exc:
        logger.warning("Could not create per-run artifact dir: %s", exc)
        return None


def _truncate(text: str, limit: int = 800) -> str:
    if len(text) <= limit:
        return text
    return text[:limit] + f"\n…[truncated {len(text) - limit} chars]"


def _dump(run_dir: Path | None, slide_index: int, name: str, content: str) -> None:
    """Best-effort artifact dump. Silent on failure — logging is the fallback."""
    if run_dir is None:
        return
    try:
        (run_dir / f"slide_{slide_index:03d}_{name}").write_text(content, encoding="utf-8")
    except OSError as exc:
        logger.debug("Could not dump %s for slide %d: %s", name, slide_index, exc)

_HTML_TEMPLATE = """\
<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=1280, height=720">
<style>
* {{ box-sizing: border-box; margin: 0; padding: 0; }}
html, body {{ width: 1280px; height: 720px; overflow: hidden; }}
{css}
.speaker-notes {{ display: none; }}
</style>
</head>
<body>
<div class="slide">
{body}
</div>
</body>
</html>
"""


class SlideBrief(BaseModel):
    """Structured input for one slide. The HTML writer sub-agent renders this.

    Keep briefs small — the writer's job is to choose layout and compose
    the course's design vocabulary; the brief just supplies the pedagogical
    content.
    """

    title: str = Field(..., description="Slide title (becomes <h1>).")
    layout: str = Field(
        default="bullets",
        description=(
            "Layout hint for the writer (advisory). One of: bullets, objectives, "
            "two-col, grid, callout, hero, code, summary. The writer picks the "
            "actual CSS composition; this just signals intent."
        ),
    )
    key_points: list[str] = Field(
        default_factory=list,
        description=(
            "Primary content for the slide. For bullets/objectives/summary "
            "layouts these become list items; for grid layouts each becomes "
            "a card title; for two-col, the first two are the column headings."
        ),
    )
    body: str = Field(
        default="",
        description=(
            "Optional richer context (1-3 short sentences) the writer can use "
            "to flesh out callouts, hero text, or card bodies."
        ),
    )
    code: str = Field(
        default="",
        description=(
            "Optional code snippet for code layouts. Keep ≤8 lines, ≤60 chars/line."
        ),
    )
    citations: list[str] = Field(
        default_factory=list,
        description=(
            "Optional citation labels (e.g. ['[^1] Apiola et al., 2022']) "
            "rendered as footnote text on the slide."
        ),
    )
    speaker_notes: str = Field(
        ...,
        description=(
            "3-5 sentences of the lecturer's script. Embedded in the PPTX "
            "notesSlide (not visible on the slide itself)."
        ),
    )


# ---------------------------------------------------------------------------
# Sub-agent infrastructure (mirrors slides_agent/tools/ModifySlide.py but
# scoped to this tool: no template registry, no critique loop, one call per
# slide).
# ---------------------------------------------------------------------------


def _read_html_writer_instructions() -> str:
    try:
        return _HTML_WRITER_INSTRUCTIONS.read_text(encoding="utf-8").strip()
    except OSError:
        return "You generate slide body HTML. Return only the HTML fragment."


def _make_html_writer_agent() -> tuple[Agent, bool]:
    """Create a fresh, stateless HTML-writer sub-agent.

    Model priority:
      1. ANTHROPIC_API_KEY → Claude Sonnet 4.6 (best HTML quality)
      2. AsyncOpenAI default (env vars) → gpt-5.3-codex
    """
    anthropic_key = os.getenv("ANTHROPIC_API_KEY")
    is_codex = False
    if anthropic_key:
        model = LitellmModel(model=_HTML_WRITER_MODEL_CLAUDE, api_key=anthropic_key)
        logger.info("HTML writer sub-agent: Claude (%s)", _HTML_WRITER_MODEL_CLAUDE)
    else:
        from agents import OpenAIResponsesModel

        client = AsyncOpenAI()
        is_codex = not str(client.base_url).startswith("https://api.openai.com")
        if is_codex:
            from dataclasses import replace

            class _CodexModel(OpenAIResponsesModel):
                async def _fetch_response(self, system_instructions, input, model_settings, *args, **kwargs):
                    model_settings = replace(model_settings, truncation=None)
                    return await super()._fetch_response(
                        system_instructions, input, model_settings, *args, **kwargs
                    )

            model = _CodexModel(model=_HTML_WRITER_MODEL_OAI, openai_client=client)
            logger.info(
                "HTML writer sub-agent: OpenAI/Codex (%s @ %s) — "
                "ANTHROPIC_API_KEY not set",
                _HTML_WRITER_MODEL_OAI, client.base_url,
            )
        else:
            model = OpenAIResponsesModel(model=_HTML_WRITER_MODEL_OAI, openai_client=client)
            logger.info(
                "HTML writer sub-agent: OpenAI (%s) — ANTHROPIC_API_KEY not set",
                _HTML_WRITER_MODEL_OAI,
            )

    instructions = _read_html_writer_instructions()
    logger.debug("HTML writer instructions loaded: %d chars", len(instructions))

    agent = Agent(
        name="Educator Slide HTML Writer",
        description="Generates one lecture slide's body HTML from a structured brief.",
        instructions=instructions,
        tools=[],
        model=model,
        model_settings=ModelSettings(
            reasoning=Reasoning(effort="medium", summary="auto"),
            verbosity="low",
            store=False if is_codex else None,
        ),
    )
    return agent, is_codex


async def _agent_get_response(agent: Agent, prompt: str, *, use_stream: bool) -> Any:
    if use_stream:
        stream = agent.get_response_stream(prompt)
        deltas: list[str] = []
        async for event in stream:
            data = getattr(event, "data", None)
            if data is not None:
                delta = getattr(data, "delta", None)
                if isinstance(delta, str):
                    deltas.append(delta)
        result = await stream.wait_final_result()
        if result is not None and not getattr(result, "final_output", None) and deltas:
            try:
                result.final_output = "".join(deltas)
            except Exception:
                pass
        return result
    return await agent.get_response(prompt)


def _build_writer_prompt(
    brief: SlideBrief,
    css_content: str,
    slide_index: int,
    total_slides: int,
) -> str:
    """Compose the per-slide prompt for the HTML-writer sub-agent."""
    parts = [
        f"You are generating slide {slide_index} of {total_slides} for a lecture deck.",
        "",
        "=== COURSE STYLE.CSS (authoritative design system) ===",
        css_content.strip(),
        "=== END STYLE.CSS ===",
        "",
        "=== SLIDE BRIEF ===",
        f"Title: {brief.title}",
        f"Layout hint: {brief.layout}",
    ]
    if brief.key_points:
        parts.append("Key points:")
        for kp in brief.key_points:
            parts.append(f"  - {kp}")
    if brief.body:
        parts.append(f"Body context: {brief.body}")
    if brief.code:
        parts.extend(["Code snippet:", "```", brief.code.rstrip(), "```"])
    if brief.citations:
        parts.append("Citations to surface on this slide:")
        for c in brief.citations:
            parts.append(f"  - {c}")
    parts.extend(
        [
            "=== END BRIEF ===",
            "",
            "Return ONLY the HTML body fragment for this slide (the content inside "
            "<div class=\"slide\">…</div>). No markdown fences, no commentary, no "
            "<html>/<body> tags, no <style> blocks, no <div class=\"speaker-notes\">.",
        ]
    )
    return "\n".join(parts)


def _extract_html_fragment(text: str) -> tuple[str, list[str]]:
    """Pull the HTML fragment out of the sub-agent's raw output.

    Handles ```html fences, ``` fences, and bare HTML. Strips any stray
    <html>/<body>/<style> the model may have emitted despite instructions.

    Returns (fragment, debug_steps) where debug_steps lists each
    transformation applied — useful for diagnosing over-aggressive
    stripping. An empty fragment means the model emitted nothing useful.
    """
    steps: list[str] = []
    raw = (text or "").strip()
    if not raw:
        steps.append("input empty")
        return "", steps
    steps.append(f"input {len(raw)} chars")

    fence = re.search(r"```(?:html)?\s*(.*?)```", raw, flags=re.IGNORECASE | re.DOTALL)
    if fence:
        raw = fence.group(1).strip()
        steps.append(f"unwrapped ``` fence → {len(raw)} chars")

    body_match = re.search(r"<body[^>]*>(.*?)</body>", raw, flags=re.IGNORECASE | re.DOTALL)
    if body_match:
        raw = body_match.group(1).strip()
        steps.append(f"peeled <body> → {len(raw)} chars")

    slide_match = re.search(
        r'<div[^>]*class=["\'][^"\']*\bslide\b[^"\']*["\'][^>]*>(.*)</div>',
        raw,
        flags=re.IGNORECASE | re.DOTALL,
    )
    if slide_match:
        raw = slide_match.group(1).strip()
        steps.append(f"peeled outer .slide wrapper → {len(raw)} chars")

    pre = len(raw)
    raw = re.sub(r"<style[^>]*>.*?</style>", "", raw, flags=re.IGNORECASE | re.DOTALL)
    if len(raw) != pre:
        steps.append(f"stripped <style> blocks → {len(raw)} chars")

    pre = len(raw)
    raw = re.sub(
        r'<div[^>]*class=["\'][^"\']*\bspeaker-notes\b[^"\']*["\'][^>]*>.*?</div>',
        "",
        raw,
        flags=re.IGNORECASE | re.DOTALL,
    )
    if len(raw) != pre:
        steps.append(f"stripped speaker-notes div → {len(raw)} chars")

    raw = raw.strip()
    if not raw:
        steps.append("WARNING: nothing left after stripping")
    return raw, steps


# ---------------------------------------------------------------------------
# Tool
# ---------------------------------------------------------------------------


class GenerateEducatorSlides(BaseTool):
    """Generate a PPTX lecture deck from structured slide briefs.

    Pass `slides`: a list of `SlideBrief` objects in presentation order.
    For each brief, this tool calls an internal HTML-writer sub-agent
    (Claude Sonnet 4.6 by default, OpenAI fallback) that turns the brief
    into a styled HTML body using the course's `style.css`. The fragments
    are wrapped, rendered to PPTX via `html2pptx_runner.js`, and written
    through the active `FileOpsBackend`. React to the response as you
    would `write_file` — see shared instructions ("Writing files").
    """

    slides: list[SlideBrief] = Field(
        ...,
        description=(
            "Structured briefs, one per slide, in presentation order. Each "
            "brief carries the slide's title, key points, optional body/code/"
            "citations, and the lecturer's speaker notes."
        ),
    )
    css_content: str = Field(
        ...,
        description="Full content of the course's style.css.",
    )
    output_path: str = Field(
        ...,
        description=(
            "Destination path for the PPTX, relative to the course root. "
            "E.g. 'lessons/L1-intro/slides.pptx'."
        ),
    )
    layout: str = Field(
        default="LAYOUT_16x9_1280",
        description="Presentation layout. Default LAYOUT_16x9_1280 (1280x720).",
    )

    async def run(self) -> str:
        if not self.slides:
            logger.warning("GenerateEducatorSlides called with empty slides list.")
            return "Error: slides is empty."
        if not RUNNER_JS.exists():
            logger.error("html2pptx_runner.js missing at %s", RUNNER_JS)
            return (
                f"Error: html2pptx_runner.js not found at {RUNNER_JS}. "
                "Ensure node_modules is installed at the repo root."
            )
        node_modules = RUNNER_JS.parent.parent.parent / "node_modules"
        if not node_modules.exists():
            logger.error("node_modules missing at %s", node_modules)
            return "Error: node_modules not found. Run 'npm install' at the repo root."

        total = len(self.slides)
        run_dir = _make_run_artifact_dir()
        logger.info(
            "GenerateEducatorSlides: %d slides → %s (layout=%s); artifacts: %s",
            total, self.output_path, self.layout, run_dir or "(disabled)",
        )
        logger.debug("css_content: %d chars", len(self.css_content))

        writer, is_codex = _make_html_writer_agent()
        html_bodies: list[str] = []
        failures: list[str] = []

        for i, brief in enumerate(self.slides, start=1):
            logger.info(
                "slide %d/%d  layout=%-9s  title=%r  key_points=%d  code=%s",
                i, total, brief.layout, brief.title,
                len(brief.key_points), "yes" if brief.code else "no",
            )
            prompt = _build_writer_prompt(brief, self.css_content, i, total)
            logger.debug("slide %d prompt: %d chars", i, len(prompt))
            _dump(run_dir, i, "prompt.txt", prompt)
            _dump(run_dir, i, "brief.json", brief.model_dump_json(indent=2))

            try:
                result = await _agent_get_response(writer, prompt, use_stream=is_codex)
            except Exception as exc:
                logger.exception("slide %d sub-agent raised: %s", i, exc)
                failures.append(f"slide {i} ({brief.title}): sub-agent error: {exc}")
                html_bodies.append(_fallback_html(brief))
                _dump(run_dir, i, "FAILURE.txt", f"sub-agent raised: {exc}")
                continue

            if result is None:
                logger.warning("slide %d: sub-agent returned None (likely API/rate-limit)", i)
                failures.append(f"slide {i} ({brief.title}): sub-agent returned None")
                html_bodies.append(_fallback_html(brief))
                _dump(run_dir, i, "FAILURE.txt", "sub-agent returned None")
                continue

            output_text = str(getattr(result, "final_output", "") or "")
            logger.debug(
                "slide %d raw output: %d chars\n%s",
                i, len(output_text), _truncate(output_text),
            )
            _dump(run_dir, i, "raw_response.txt", output_text)

            fragment, extraction_steps = _extract_html_fragment(output_text)
            for step in extraction_steps:
                logger.debug("slide %d extract: %s", i, step)

            if not fragment:
                logger.warning(
                    "slide %d: extraction produced empty fragment — falling back. "
                    "raw output starts with: %r",
                    i, output_text[:200],
                )
                failures.append(f"slide {i} ({brief.title}): empty sub-agent output")
                html_bodies.append(_fallback_html(brief))
                _dump(run_dir, i, "FAILURE.txt", "extraction produced empty fragment")
                continue

            logger.info("slide %d ok: %d chars of HTML", i, len(fragment))
            _dump(run_dir, i, "fragment.html", fragment)
            html_bodies.append(fragment)

        successes = total - len(failures)
        logger.info(
            "sub-agent loop done: %d ok, %d fell back to plain bullets",
            successes, len(failures),
        )
        if failures:
            failure_note = "\n".join(f"  - {f}" for f in failures)
        else:
            failure_note = ""

        with tempfile.TemporaryDirectory(prefix="educator_slides_") as tmp:
            tmp_path = Path(tmp)
            html_files = self._write_html_files(tmp_path, html_bodies)
            # Snapshot the final per-slide HTML (post-wrapping) for inspection.
            for idx, html_file in enumerate(html_files, start=1):
                _dump(run_dir, idx, "final_wrapped.html", Path(html_file).read_text(encoding="utf-8"))
            pptx_path = tmp_path / "slides.pptx"
            logger.debug("running html2pptx_runner.js on %d files", len(html_files))
            error = self._run_converter(html_files, pptx_path)
            if error:
                logger.error("html2pptx_runner.js failed: %s", error)
                return error
            try:
                pptx_bytes = pptx_path.read_bytes()
            except OSError as exc:
                logger.exception("could not read generated PPTX")
                return f"Failed to read generated PPTX: {exc}"
            logger.debug("PPTX generated: %d bytes", len(pptx_bytes))
            pptx_bytes = self._add_speaker_notes(pptx_bytes)

        backend = get_backend()
        outcome = backend.write_file(Path(self.output_path), pptx_bytes)
        logger.info("write outcome: %s", type(outcome).__name__)

        if isinstance(outcome, Pending):
            msg = (
                f"Slide deck generated ({total} slides) — pending approval.\n"
                f"proposal_id: {outcome.proposal_id}\n"
                f"path: {outcome.path}\n\n"
                f"{outcome.diff}\n"
                f"Reply `/approve {outcome.proposal_id}` to save, or "
                f"`/reject {outcome.proposal_id} <feedback>` to discard."
            )
        elif isinstance(outcome, Accepted):
            msg = f"Slide deck saved: {outcome.path} ({total} slides)"
        elif isinstance(outcome, Rejected):
            msg = f"Proposal rejected: {outcome.reason}"
        elif isinstance(outcome, Failed):
            msg = f"Write failed: {outcome.error}"
        else:
            msg = f"Unexpected outcome: {outcome!r}"

        if failure_note:
            msg += (
                f"\n\nNote: {len(failures)} slide(s) used a plain fallback layout "
                f"because the HTML writer failed:\n{failure_note}"
            )
        return msg

    def _write_html_files(self, tmp_path: Path, html_bodies: list[str]) -> list[str]:
        paths = []
        for i, (brief, body) in enumerate(zip(self.slides, html_bodies)):
            full_body = body
            if brief.speaker_notes.strip():
                # Escape only </div> sequences that would close the wrapper early.
                notes = brief.speaker_notes.strip().replace("</div>", "&lt;/div&gt;")
                full_body = f'{body}\n<div class="speaker-notes">{notes}</div>'
            html = _HTML_TEMPLATE.format(css=self.css_content, body=full_body)
            slide_file = tmp_path / f"slide_{i + 1:03d}.html"
            slide_file.write_text(html, encoding="utf-8")
            paths.append(str(slide_file))
        return paths

    def _run_converter(self, html_files: list[str], output_path: Path) -> str | None:
        cmd = [
            "node",
            str(RUNNER_JS),
            "--output", str(output_path),
            "--layout", self.layout,
            "--",
            *html_files,
        ]
        try:
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=300,
                cwd=str(RUNNER_JS.parent.parent.parent),
            )
        except subprocess.TimeoutExpired:
            return "Error: PPTX conversion timed out after 300 s."
        except FileNotFoundError:
            return "Error: Node.js not found. Please install Node.js."

        if result.returncode != 0:
            err = result.stderr.strip() or result.stdout.strip()
            return f"Error converting HTML to PPTX:\n{err}"
        return None

    def _add_speaker_notes(self, pptx_bytes: bytes) -> bytes:
        """Embed speaker notes from each brief into the PPTX notesSlide."""
        from io import BytesIO

        from pptx import Presentation

        prs = Presentation(BytesIO(pptx_bytes))
        for i, slide in enumerate(prs.slides):
            if i >= len(self.slides):
                continue
            notes = self.slides[i].speaker_notes.strip()
            if not notes:
                continue
            slide.notes_slide.notes_text_frame.text = notes

        out = BytesIO()
        prs.save(out)
        return out.getvalue()


def _fallback_html(brief: SlideBrief) -> str:
    """Plain bullets layout when the sub-agent fails on a slide.

    The deck still ships; the user gets a clear note about which slides
    fell back so they can re-prompt for a regeneration.
    """
    parts = [f"<h1>{_escape(brief.title)}</h1>"]
    if brief.key_points:
        parts.append("<ul>")
        for kp in brief.key_points:
            parts.append(f"  <li>{_escape(kp)}</li>")
        parts.append("</ul>")
    if brief.body:
        parts.append(f"<p>{_escape(brief.body)}</p>")
    if brief.code:
        parts.append(f"<pre><code>{_escape(brief.code)}</code></pre>")
    if brief.citations:
        for c in brief.citations:
            parts.append(f'<div class="footnote">{_escape(c)}</div>')
    return "\n".join(parts)


def _escape(s: str) -> str:
    return (
        s.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
    )
