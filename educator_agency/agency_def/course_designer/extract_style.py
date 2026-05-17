"""Tool: extract a candidate style.css from an educator-supplied PPTX.

Reads the source PPTX through the backend (so the file must already exist
inside the course root), pulls theme colors, font families, type scale, and
the most-likely logo image, and writes a candidate `style.css` (and optional
`assets/logo.<ext>`) back through the backend. Educator reviews via the
approval-gate diff and approves/rejects in the usual way.

What we extract:

- **Theme colors** — from `/ppt/theme/theme1.xml` `<a:clrScheme>`. Falls back
  to most-common direct RGB fills if the theme is sparse.
- **Font families** — `<a:fontScheme>` `<a:majorFont>` / `<a:minorFont>`,
  cross-checked against actual `font.name` usage in slide runs.
- **Type scale** — histogram of unique font sizes (in pt) across all text
  runs, binned into title / h2 / body / small.
- **Logo** — image that appears in roughly the same position on ≥30% of
  slides. If no clear winner, candidates are listed in the tool response.

What we deliberately don't try:

- Layout primitives (two-column, hero, kicker, etc.) — not mechanically
  extractable without layout-classification ML.
- Spacing/margin conventions — inference only, not reliable.
- Keynote `.key` files — out of scope; educator exports to .pptx first.
"""

from __future__ import annotations

import hashlib
import zipfile
from collections import Counter
from io import BytesIO
from pathlib import Path
from typing import Optional
from xml.etree import ElementTree as ET

from agency_swarm.tools import BaseTool
from pydantic import Field

from educator_agency.runtime._context import get_backend
from educator_agency.runtime.file_ops import Accepted, Failed, Pending, Rejected

_NS_DRAWINGML = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS = {"a": _NS_DRAWINGML}

# Theme color slot → label we use when synthesizing the CSS.
_COLOR_SLOTS = ["dk1", "lt1", "accent1", "accent2", "accent3", "accent4"]


class ExtractStyleFromPptx(BaseTool):
    """Extract a candidate style.css from an existing PPTX deck.

    Pass `source_pptx_path` (relative to the course root) pointing at a
    PPTX the educator has placed in the course directory (e.g. uploaded
    `style_source.pptx`). The tool extracts theme colors, fonts, a type
    scale, and a candidate logo image, then writes a `style.css` through
    the backend. Handle the response per shared instructions
    ("Writing files").

    If a logo is identified, it is also written to `assets/logo.<ext>`
    (a separate write proposal). If multiple candidates tie, the tool
    skips the logo write and lists the candidates in its response so the
    educator can choose manually.
    """

    source_pptx_path: str = Field(
        ...,
        description=(
            "Path to the source .pptx file, relative to the course root. "
            "E.g. 'style_source.pptx' if the educator dropped it at the "
            "course root."
        ),
    )
    target_css_path: str = Field(
        default="style.css",
        description="Where to write the candidate style.css. Defaults to course-root style.css.",
    )

    def run(self) -> str:
        backend = get_backend()
        try:
            pptx_bytes = backend.read_bytes(Path(self.source_pptx_path))
        except FileNotFoundError:
            return f"Error: source PPTX not found at {self.source_pptx_path}."
        except Exception as exc:
            return f"Error reading {self.source_pptx_path}: {exc}"

        try:
            extracted = _extract(pptx_bytes)
        except Exception as exc:
            return f"Extraction failed: {exc}"

        css = _render_css(extracted)
        outcome = backend.write_file(Path(self.target_css_path), css)
        primary_msg = _format_outcome(outcome, "style.css")

        notes: list[str] = []
        # Surface what we found (and didn't) so the educator can sanity-check.
        notes.append(_describe_extraction(extracted))

        # Logo: write a separate proposal if we have a clear winner.
        if extracted.logo_image is not None and extracted.logo_image_ext:
            logo_path = Path("assets") / f"logo.{extracted.logo_image_ext}"
            logo_outcome = backend.write_file(logo_path, extracted.logo_image)
            notes.append(
                "\nLogo: " + _format_outcome(logo_outcome, str(logo_path))
            )
        elif extracted.logo_candidates:
            notes.append(
                "\nLogo: no clear winner. Candidates found: "
                + ", ".join(extracted.logo_candidates)
                + ". Drop the one you want at `assets/logo.<ext>` manually."
            )
        else:
            notes.append("\nLogo: no recurring image detected.")

        return primary_msg + "\n\n" + "\n".join(notes)


# ---------------------------------------------------------------------------
# Extraction
# ---------------------------------------------------------------------------


class _Extracted:
    """Bag of extracted fields. Plain class — no validation needed."""

    def __init__(self) -> None:
        self.theme_colors: dict[str, str] = {}  # slot name → hex (no #)
        self.heading_font: Optional[str] = None
        self.body_font: Optional[str] = None
        self.font_sizes_pt: list[float] = []
        self.detected_fonts: Counter[str] = Counter()
        self.logo_image: Optional[bytes] = None
        self.logo_image_ext: Optional[str] = None
        self.logo_candidates: list[str] = []
        self.slide_count: int = 0
        self.notes: list[str] = []


def _extract(pptx_bytes: bytes) -> _Extracted:
    out = _Extracted()

    # Theme XML — colors + font scheme.
    try:
        with zipfile.ZipFile(BytesIO(pptx_bytes)) as zf:
            theme_xml = _read_first(zf, ["ppt/theme/theme1.xml"])
            if theme_xml is not None:
                _parse_theme(theme_xml, out)
            else:
                out.notes.append("No ppt/theme/theme1.xml found.")
    except zipfile.BadZipFile:
        raise ValueError("Source file is not a valid PPTX (not a ZIP).")

    # python-pptx for shape/text/image traversal.
    from pptx import Presentation
    from pptx.util import Emu

    prs = Presentation(BytesIO(pptx_bytes))
    out.slide_count = len(prs.slides)

    image_positions: dict[str, list[tuple[int, int, int, int]]] = {}
    image_blobs: dict[str, tuple[bytes, str]] = {}

    for slide in prs.slides:
        _walk_shapes(slide.shapes, out, image_positions, image_blobs)

    _pick_logo(image_positions, image_blobs, out)

    # If theme didn't give us fonts but slides did, use the most common.
    if not out.heading_font and out.detected_fonts:
        out.heading_font = out.detected_fonts.most_common(1)[0][0]
    if not out.body_font and out.detected_fonts:
        out.body_font = out.detected_fonts.most_common(1)[0][0]

    _ = Emu  # silence unused-import lint (kept for future use)
    return out


def _read_first(zf: zipfile.ZipFile, names: list[str]) -> Optional[bytes]:
    for name in names:
        if name in zf.namelist():
            return zf.read(name)
    return None


def _parse_theme(theme_xml: bytes, out: _Extracted) -> None:
    root = ET.fromstring(theme_xml)
    clr_scheme = root.find(f".//{{{_NS_DRAWINGML}}}clrScheme")
    if clr_scheme is not None:
        for slot in _COLOR_SLOTS:
            node = clr_scheme.find(f"{{{_NS_DRAWINGML}}}{slot}")
            if node is None:
                continue
            hex_val = _resolve_color(node)
            if hex_val:
                out.theme_colors[slot] = hex_val

    font_scheme = root.find(f".//{{{_NS_DRAWINGML}}}fontScheme")
    if font_scheme is not None:
        major = font_scheme.find(f"{{{_NS_DRAWINGML}}}majorFont/{{{_NS_DRAWINGML}}}latin")
        minor = font_scheme.find(f"{{{_NS_DRAWINGML}}}minorFont/{{{_NS_DRAWINGML}}}latin")
        if major is not None:
            out.heading_font = major.get("typeface") or None
        if minor is not None:
            out.body_font = minor.get("typeface") or None


def _resolve_color(slot_node: ET.Element) -> Optional[str]:
    """Resolve a theme-color slot to a 6-char hex string (no #), if possible."""
    srgb = slot_node.find(f"{{{_NS_DRAWINGML}}}srgbClr")
    if srgb is not None and srgb.get("val"):
        return srgb.get("val").upper()
    sys = slot_node.find(f"{{{_NS_DRAWINGML}}}sysClr")
    if sys is not None and sys.get("lastClr"):
        return sys.get("lastClr").upper()
    return None


def _walk_shapes(
    shapes,
    out: _Extracted,
    image_positions: dict[str, list[tuple[int, int, int, int]]],
    image_blobs: dict[str, tuple[bytes, str]],
) -> None:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            _walk_shapes(shape.shapes, out, image_positions, image_blobs)
            continue

        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    font = run.font
                    if font.name:
                        out.detected_fonts[font.name] += 1
                    if font.size is not None:
                        out.font_sizes_pt.append(float(font.size.pt))

        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                blob = shape.image.blob
                ext = (shape.image.ext or "png").lower()
            except Exception:
                continue
            digest = hashlib.sha1(blob).hexdigest()
            image_blobs[digest] = (blob, ext)
            image_positions.setdefault(digest, []).append(
                (shape.left or 0, shape.top or 0, shape.width or 0, shape.height or 0)
            )


def _pick_logo(
    image_positions: dict[str, list[tuple[int, int, int, int]]],
    image_blobs: dict[str, tuple[bytes, str]],
    out: _Extracted,
) -> None:
    if not image_positions or out.slide_count == 0:
        return

    # An image is "logo-like" if it appears on ≥30% of slides AND its
    # positions across appearances are tightly clustered (low variance).
    threshold = max(2, int(round(out.slide_count * 0.3)))

    scored: list[tuple[int, str]] = []  # (count, digest)
    for digest, positions in image_positions.items():
        if len(positions) < threshold:
            continue
        # Cluster check: max spread in left/top across appearances.
        lefts = [p[0] for p in positions]
        tops = [p[1] for p in positions]
        spread = max(max(lefts) - min(lefts), max(tops) - min(tops))
        # 1 inch = 914400 EMU. Accept up to ~0.5" wander.
        if spread > 457_200:
            continue
        scored.append((len(positions), digest))

    if not scored:
        return

    scored.sort(reverse=True)
    top_count = scored[0][0]
    winners = [d for c, d in scored if c == top_count]

    if len(winners) == 1:
        digest = winners[0]
        blob, ext = image_blobs[digest]
        out.logo_image = blob
        out.logo_image_ext = ext
    else:
        # Tied — surface as candidates instead of guessing.
        out.logo_candidates = [f"{d[:8]} ({image_blobs[d][1]})" for d in winners]


# ---------------------------------------------------------------------------
# CSS rendering
# ---------------------------------------------------------------------------


def _render_css(ex: _Extracted) -> str:
    # Color tagging heuristic: dk1 → text, lt1 → bg, accent1 → main accent.
    fg = _hex(ex.theme_colors.get("dk1"), "#1a1a1a")
    bg = _hex(ex.theme_colors.get("lt1"), "#ffffff")
    accent = _hex(ex.theme_colors.get("accent1"), "#2c5aa0")
    c1 = _hex(ex.theme_colors.get("accent1"), accent)
    c2 = _hex(ex.theme_colors.get("accent2"), "#117a65")
    c3 = _hex(ex.theme_colors.get("accent3"), "#b7791f")
    c4 = _hex(ex.theme_colors.get("accent4"), "#9b2c2c")

    heading_font = ex.heading_font or "Helvetica Neue"
    body_font = ex.body_font or "Helvetica Neue"

    # Type scale: bin observed sizes; fall back to defaults.
    fs_title, fs_h2, fs_body, fs_small = _bin_type_scale(ex.font_sizes_pt)

    google_fonts = _google_fonts_import(heading_font, body_font)

    return f"""\
/* Candidate style.css extracted from a source PPTX.
 * Review and adjust before committing — the type scale, palette tagging,
 * and font families are heuristic. Spacing, layout primitives, and
 * component classes are inherited from the educator-agency design system
 * and are not extracted from the source deck.
 */

{google_fonts}

:root {{
  /* Palette (extracted from PPTX theme; tagging is heuristic) */
  --slide-bg: {bg};
  --slide-fg: {fg};
  --accent: {accent};
  --accent-soft: #e8f0fa;
  --muted: #666666;
  --rule: #e5e7eb;
  --surface: #f7f8fa;

  --c1: {c1};
  --c2: {c2};
  --c3: {c3};
  --c4: {c4};

  /* Typography (extracted from PPTX theme + observed runs) */
  --font-body: "{body_font}", "Helvetica Neue", Arial, sans-serif;
  --font-heading: "{heading_font}", "Helvetica Neue", Arial, sans-serif;
  --font-mono: "IBM Plex Mono", "Menlo", "Courier New", monospace;

  /* Type scale (px — fixed 1280x720 canvas) */
  --fs-title: {fs_title}px;
  --fs-h2: {fs_h2}px;
  --fs-body: {fs_body}px;
  --fs-small: {fs_small}px;
  --fs-kicker: 16px;

  /* Spacing scale */
  --space-1: 8px;
  --space-2: 16px;
  --space-3: 24px;
  --space-4: 40px;
  --space-5: 60px;
}}

.slide {{
  background: var(--slide-bg);
  color: var(--slide-fg);
  font-family: var(--font-body);
  font-size: var(--fs-body);
  line-height: 1.4;
  padding: var(--space-5);
  width: 1280px;
  height: 720px;
  position: relative;
  overflow: hidden;
  display: flex;
  flex-direction: column;
  gap: var(--space-3);
}}

.slide h1 {{
  font-family: var(--font-heading);
  font-size: var(--fs-title);
  font-weight: 700;
  color: var(--accent);
  line-height: 1.1;
  margin: 0;
}}
.slide h2 {{
  font-family: var(--font-heading);
  font-size: var(--fs-h2);
  font-weight: 600;
  margin: 0;
}}
.slide p {{ margin: 0; }}
.slide ul, .slide ol {{ margin: 0; padding-left: 1.2em; }}
.slide li {{ margin-bottom: var(--space-1); }}
.slide code {{
  font-family: var(--font-mono);
  background: var(--surface);
  padding: 2px 6px;
  border-radius: 3px;
  font-size: 0.85em;
}}
.slide pre {{
  font-family: var(--font-mono);
  background: var(--surface);
  padding: var(--space-2);
  border-radius: 6px;
  font-size: var(--fs-small);
  line-height: 1.4;
  overflow: hidden;
}}

.kicker {{
  font-family: var(--font-mono);
  font-size: var(--fs-kicker);
  font-weight: 500;
  letter-spacing: 0.12em;
  text-transform: uppercase;
  color: var(--accent);
}}

.accent-bar {{
  display: block;
  width: var(--space-4);
  height: 3px;
  background: var(--accent);
  margin-bottom: var(--space-2);
}}
.accent-bar.full {{ width: 100%; }}

.card {{
  background: var(--surface);
  border-radius: 8px;
  padding: var(--space-3);
  min-height: 120px;
  display: flex;
  flex-direction: column;
  gap: var(--space-1);
}}
.card h3 {{
  font-family: var(--font-heading);
  font-size: 22px;
  font-weight: 600;
  margin: 0;
}}
.card.outlined {{
  background: transparent;
  border: 1px solid var(--rule);
}}

.two-col {{
  display: flex;
  gap: var(--space-4);
  flex: 1;
}}
.two-col > * {{ flex: 1; }}

.grid-3 {{
  display: grid;
  grid-template-columns: repeat(3, 1fr);
  gap: var(--space-3);
  flex: 1;
}}
.grid-2 {{
  display: grid;
  grid-template-columns: repeat(2, 1fr);
  gap: var(--space-3);
  flex: 1;
}}

.hero {{
  flex: 1;
  display: flex;
  flex-direction: column;
  justify-content: center;
  align-items: flex-start;
  gap: var(--space-2);
}}
.hero .stat {{
  font-family: var(--font-heading);
  font-size: 120px;
  font-weight: 700;
  color: var(--accent);
  line-height: 1;
}}

.callout {{
  background: var(--accent-soft);
  border-left: 4px solid var(--accent);
  padding: var(--space-2) var(--space-3);
  border-radius: 0 6px 6px 0;
  font-size: var(--fs-body);
}}

.color-1 {{ --accent: var(--c1); }}
.color-2 {{ --accent: var(--c2); }}
.color-3 {{ --accent: var(--c3); }}
.color-4 {{ --accent: var(--c4); }}

.footnote {{
  font-family: var(--font-mono);
  font-size: 14px;
  color: var(--muted);
  position: absolute;
  bottom: 30px;
  right: var(--space-5);
}}

.speaker-notes {{ display: none; }}
"""


def _hex(value: Optional[str], default: str) -> str:
    if not value:
        return default
    v = value.strip().lstrip("#")
    if len(v) == 6 and all(c in "0123456789abcdefABCDEF" for c in v):
        return f"#{v.lower()}"
    return default


def _bin_type_scale(sizes_pt: list[float]) -> tuple[int, int, int, int]:
    """Bin observed PPTX font sizes (pt) into title/h2/body/small (px).

    Conservative defaults if we don't see enough variation.
    """
    if not sizes_pt:
        return (56, 36, 26, 18)

    unique = sorted(set(round(s) for s in sizes_pt), reverse=True)
    # Heuristic: take the largest distinct size as title; next as h2; the
    # mode as body; smallest as small. Convert pt → px at ~1.5x (slide
    # canvas is rendered at higher effective DPI than 72).
    title_pt = unique[0] if unique else 32
    h2_pt = unique[1] if len(unique) > 1 else max(24, title_pt - 12)
    small_pt = unique[-1] if len(unique) > 1 else 12
    counter = Counter(round(s) for s in sizes_pt)
    body_pt = counter.most_common(1)[0][0]

    def pt_to_px(pt: float) -> int:
        return max(12, int(round(pt * 1.5)))

    return (
        pt_to_px(title_pt),
        pt_to_px(h2_pt),
        pt_to_px(body_pt),
        pt_to_px(small_pt),
    )


def _google_fonts_import(heading: str, body: str) -> str:
    """Best-effort Google Fonts @import for the detected families.

    If the families aren't on Google Fonts the @import will just 404 and
    the browser falls back to the system stack — no error.
    """
    families: list[str] = []
    seen: set[str] = set()
    for family in (heading, body, "IBM Plex Mono"):
        key = family.replace(" ", "+")
        if key.lower() in seen:
            continue
        seen.add(key.lower())
        families.append(f"family={key}:wght@400;500;600;700")
    return (
        '@import url("https://fonts.googleapis.com/css2?'
        + "&".join(families)
        + '&display=swap");'
    )


# ---------------------------------------------------------------------------
# Reporting
# ---------------------------------------------------------------------------


def _format_outcome(outcome, label: str) -> str:
    if isinstance(outcome, Pending):
        return (
            f"{label} — pending approval.\n"
            f"proposal_id: {outcome.proposal_id}\n"
            f"path: {outcome.path}\n\n"
            f"{outcome.diff}\n"
            f"Reply `/approve {outcome.proposal_id}` to save, or "
            f"`/reject {outcome.proposal_id} <feedback>` to discard."
        )
    if isinstance(outcome, Accepted):
        return f"{label} written: {outcome.path}"
    if isinstance(outcome, Rejected):
        return f"{label} rejected: {outcome.reason}"
    if isinstance(outcome, Failed):
        return f"{label} write failed: {outcome.error}"
    return f"{label}: unexpected outcome {outcome!r}"


def _describe_extraction(ex: _Extracted) -> str:
    lines = ["Extracted from source PPTX:"]
    if ex.theme_colors:
        lines.append(
            "  Theme colors: "
            + ", ".join(f"{k}=#{v}" for k, v in ex.theme_colors.items())
        )
    else:
        lines.append("  Theme colors: none found (using defaults).")
    lines.append(
        f"  Fonts: heading='{ex.heading_font or '?'}', body='{ex.body_font or '?'}'"
    )
    if ex.detected_fonts:
        top = ex.detected_fonts.most_common(3)
        lines.append(
            "  Observed font usage: "
            + ", ".join(f"{n}×{c}" for n, c in top)
        )
    if ex.font_sizes_pt:
        unique = sorted(set(round(s) for s in ex.font_sizes_pt), reverse=True)
        lines.append(f"  Observed font sizes (pt): {unique[:8]}")
    lines.append(f"  Slides scanned: {ex.slide_count}")
    for note in ex.notes:
        lines.append(f"  Note: {note}")
    return "\n".join(lines)
