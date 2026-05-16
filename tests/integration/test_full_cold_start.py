"""Full §6.1 cold-start integration test (LLM-driven).

Exercises the complete PLAN.md → SlidesAgent → ApprovalBuffer → /approve →
slides.pptx pipeline with a real LLM call. Skipped unless:
  - OPENAI_API_KEY is set in the environment
  - Node.js and node_modules are present (for html2pptx_runner.js)

Run explicitly with:
  pytest -m llm tests/integration/test_full_cold_start.py -v

Per `.claude/plans/i-want-to-start-stateless-oasis.md` §6.1, §8.
"""

from __future__ import annotations

import os
from io import BytesIO
from pathlib import Path

import pytest

from educator_agency.runtime._context import set_backend
from educator_agency.runtime.file_ops import ApprovalBuffer, ApprovalGatingBackend, LocalFsBackend
from educator_agency.runtime.slash_commands import DirectResponse, dispatch


def _has_api_key() -> bool:
    return bool(os.getenv("OPENAI_API_KEY") or os.getenv("ANTHROPIC_API_KEY"))


pytestmark = pytest.mark.llm


@pytest.fixture
def live_backend(course_root: Path):
    """Wire a real ApprovalGatingBackend for LLM tests."""
    buf = ApprovalBuffer()
    inner = LocalFsBackend(course_root)
    gating = ApprovalGatingBackend(inner, buf)
    set_backend(gating)
    return gating, buf, course_root


@pytest.fixture
def educator_agency(live_backend):
    """Instantiate the educator-agency in dry-run mode; skip if no API key."""
    if not _has_api_key():
        pytest.skip("No API key available (set OPENAI_API_KEY to run LLM tests)")
    from educator_agency.agency_def.agency import create_educator_agency
    return create_educator_agency()


@pytest.mark.usefixtures("require_node")
def test_slides_generated_and_approved_via_llm(educator_agency, live_backend):
    """SlidesAgent reads PLAN.md and generates a slide proposal via the LLM.

    This test exercises the real agency without a mock. Expected outcome:
      1. Orchestrator delegates to SlidesAgent via SendMessage
      2. SlidesAgent reads PLAN.md, PEDAGOGY.md, style.css
      3. SlidesAgent calls generate_educator_slides → Pending proposal
      4. ApprovalBuffer has exactly one proposal
      5. /approve writes slides.pptx to disk
      6. slides.pptx opens and contains ≥ 3 slides with ≥ 1 notes slide
    """
    from agents import RunConfig

    gating, buf, course_root = live_backend

    # Use a cheaper model for the integration test (cost ceiling per plan §6.11)
    run_config = RunConfig(model="gpt-4o-mini")

    result = educator_agency.get_response_sync(
        message=(
            "Generate the slides for lesson 1 (Introduction to hashing). "
            "Read the PLAN.md at lessons/L1-introduction-to-hashing/PLAN.md. "
            "Produce a lecture-style deck and present the proposal for approval."
        ),
        run_config=run_config,
    )

    assert buf, (
        "Expected at least one pending proposal in the ApprovalBuffer after the "
        f"agency run. Agency output: {result}"
    )

    # Approve the first proposal (there should be exactly one for slides.pptx)
    proposal = next(buf.pending())
    slash_result = dispatch(f"/approve {proposal.proposal_id}", gating)
    assert isinstance(slash_result, DirectResponse)
    assert "accepted" in slash_result.text.lower(), slash_result.text

    pptx_path = course_root / "lessons" / "L1-introduction-to-hashing" / "slides.pptx"
    assert pptx_path.exists(), "slides.pptx not written after approval"
    assert pptx_path.stat().st_size > 5000, "slides.pptx suspiciously small"

    from pptx import Presentation

    prs = Presentation(BytesIO(pptx_path.read_bytes()))
    assert len(prs.slides) >= 3, f"Expected ≥3 slides, got {len(prs.slides)}"

    notes_slides = [s for s in prs.slides if s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip()]
    assert len(notes_slides) >= 1, "Expected at least 1 slide with embedded speaker notes"
