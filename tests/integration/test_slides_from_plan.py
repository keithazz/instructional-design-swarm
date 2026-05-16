"""Integration test: hand-written PLAN.md → slides.pptx.

This is the §6.1 "carved slice" integration test from the implementation plan.
It exercises the full pipeline WITHOUT an LLM call:

  FileOpsBackend → GenerateEducatorSlides tool → html2pptx_runner.js →
  python-pptx speaker-notes pass → ApprovalBuffer → /approve → file on disk

This proves the architectural seam works: the tool generates PPTX, routes it
through the approval gate as Pending, and the slash-command handler commits it.

Requires: Node.js, node_modules (html2pptx_runner.js, playwright)
Skipped automatically if those are absent.

LLM-driven tests (SlidesAgent generating HTML from PLAN.md) are in
`test_slides_llm.py` and marked @pytest.mark.llm.

Per `.claude/plans/i-want-to-start-stateless-oasis.md` §6.1, §8.
"""

from __future__ import annotations

from io import BytesIO
from pathlib import Path

import pytest

from educator_agency.agency_def.slides_agent.generate_slides import GenerateEducatorSlides
from educator_agency.runtime.file_ops import Accepted, Pending
from educator_agency.runtime.slash_commands import DirectResponse, dispatch
from tests.integration.conftest import SAMPLE_STYLE_CSS

MINIMAL_SLIDES_HTML = [
    # Slide 1: title slide
    """\
<h1>Introduction to Cryptographic Hashing</h1>
<ul>
  <li>LO-1.1: Define a cryptographic hash function</li>
  <li>LO-1.2: Three security properties (informally)</li>
  <li>LO-1.3: Real-world application categories</li>
</ul>
<div class="speaker-notes">
Welcome to lesson 1. Today we introduce cryptographic hash functions. We will
start with a motivating example from everyday password authentication and build
up to a precise characterisation of what makes a hash function cryptographic.
By the end you should be able to distinguish cryptographic from non-cryptographic
hashes and name the three core security properties.
</div>
""",
    # Slide 2: motivating example
    """\
<h1>Why Hash Functions Matter</h1>
<h2>Password authentication example</h2>
<ul>
  <li>Server stores <code>H(password)</code>, not the password itself</li>
  <li>On login: compute <code>H(input)</code> and compare</li>
  <li>Attacker who steals the database cannot reverse <code>H</code></li>
</ul>
<div class="footnote">[^1] Menezes et al., 1996</div>
<div class="speaker-notes">
Consider what happens when a password database is breached. If passwords were
stored in plaintext, every account is immediately compromised. Instead, we store
the digest — the output of a hash function applied to the password. Even with the
full database, an attacker cannot recover the original passwords without inverting
the hash, which the function is designed to make infeasible.
</div>
""",
    # Slide 3: definition
    """\
<h1>Definition</h1>
<p>A <strong>cryptographic hash function</strong> is a deterministic function:</p>
<p><code>H : {0,1}* → {0,1}^n</code></p>
<ul>
  <li>Efficiently computable</li>
  <li>Fixed-length output regardless of input length</li>
  <li>Designed to resist computationally-bounded adversaries</li>
</ul>
<div class="speaker-notes">
The formal definition. H maps bitstrings of arbitrary length to bitstrings of
fixed length n — for SHA-256, n is 256 bits. Efficiently computable means a
modern CPU computes it in microseconds. The adversary model is the key
distinction from non-cryptographic hashes: we assume an attacker with
polynomial-time compute and ask whether they can find structure in the mapping.
</div>
""",
]


@pytest.mark.usefixtures("require_node")
def test_generate_slides_full_pipeline(wired_backend):
    """Full pipeline: tool generates PPTX, proposal created, /approve commits file."""
    gating, buf, course_root = wired_backend

    tool = GenerateEducatorSlides(
        slides_html=MINIMAL_SLIDES_HTML,
        css_content=SAMPLE_STYLE_CSS,
        output_path="lessons/L1-introduction-to-hashing/slides.pptx",
    )
    result = tool.run()

    assert "proposal_id" in result, f"Expected proposal_id in result, got: {result[:200]}"
    assert len(buf) == 1, "Expected exactly one pending proposal"

    proposal = next(buf.pending())
    proposal_id = proposal.proposal_id

    slash_result = dispatch(f"/approve {proposal_id}", gating)
    assert isinstance(slash_result, DirectResponse)
    assert "accepted" in slash_result.text.lower()

    pptx_path = course_root / "lessons" / "L1-introduction-to-hashing" / "slides.pptx"
    assert pptx_path.exists(), "slides.pptx not written after approval"
    assert pptx_path.stat().st_size > 1000, "slides.pptx suspiciously small"

    from pptx import Presentation

    prs = Presentation(BytesIO(pptx_path.read_bytes()))
    assert len(prs.slides) >= 3, (
        f"Expected ≥3 slides but got {len(prs.slides)}"
    )


@pytest.mark.usefixtures("require_node")
def test_rejected_proposal_does_not_write(wired_backend):
    """Rejecting a proposal leaves slides.pptx absent from disk."""
    gating, buf, course_root = wired_backend

    tool = GenerateEducatorSlides(
        slides_html=MINIMAL_SLIDES_HTML,
        css_content=SAMPLE_STYLE_CSS,
        output_path="lessons/L1-introduction-to-hashing/slides.pptx",
    )
    result = tool.run()
    assert "proposal_id" in result

    proposal_id = next(buf.pending()).proposal_id
    slash_result = dispatch(f"/reject {proposal_id} too many slides on one topic", gating)
    assert isinstance(slash_result, DirectResponse)
    assert "rejected" in slash_result.text.lower()

    pptx_path = course_root / "lessons" / "L1-introduction-to-hashing" / "slides.pptx"
    assert not pptx_path.exists(), "slides.pptx should not exist after rejection"
    assert len(buf) == 0, "Buffer should be empty after rejection"


@pytest.mark.usefixtures("require_node")
def test_speaker_notes_embedded_in_pptx(wired_backend):
    """Speaker notes from <div class='speaker-notes'> appear in the PPTX notes slides."""
    gating, buf, course_root = wired_backend

    tool = GenerateEducatorSlides(
        slides_html=MINIMAL_SLIDES_HTML,
        css_content=SAMPLE_STYLE_CSS,
        output_path="lessons/L1-introduction-to-hashing/slides.pptx",
    )
    tool.run()

    proposal_id = next(buf.pending()).proposal_id
    dispatch(f"/approve {proposal_id}", gating)

    pptx_path = course_root / "lessons" / "L1-introduction-to-hashing" / "slides.pptx"
    from pptx import Presentation

    prs = Presentation(BytesIO(pptx_path.read_bytes()))
    notes_texts = [
        slide.notes_slide.notes_text_frame.text.strip()
        for slide in prs.slides
        if slide.has_notes_slide
    ]
    # At least the first two slides have non-empty notes.
    non_empty = [t for t in notes_texts if t]
    assert len(non_empty) >= 2, (
        f"Expected ≥2 slides with speaker notes, got {len(non_empty)}: {notes_texts}"
    )
    assert "cryptographic" in notes_texts[0].lower()


@pytest.mark.usefixtures("require_node")
def test_chaos_backend_reports_rejection_without_writing(course_root, tmp_path):
    """ChaosBackend returns Rejected — file never written."""
    from educator_agency.runtime.file_ops import ChaosBackend, LocalFsBackend
    from educator_agency.runtime._context import set_backend

    inner = LocalFsBackend(course_root)
    chaos = ChaosBackend(inner, reject_reason="chaos: always reject")
    set_backend(chaos)

    tool = GenerateEducatorSlides(
        slides_html=MINIMAL_SLIDES_HTML,
        css_content=SAMPLE_STYLE_CSS,
        output_path="lessons/L1-introduction-to-hashing/slides.pptx",
    )
    result = tool.run()

    assert "rejected" in result.lower(), f"Expected rejection message, got: {result[:200]}"
    pptx_path = course_root / "lessons" / "L1-introduction-to-hashing" / "slides.pptx"
    assert not pptx_path.exists()
